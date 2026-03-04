"""
generate_ppt.py
Generates a professional PowerPoint presentation for
Waste-Not: Circular Economy Tracker (Samsung Hackathon)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour Palette ────────────────────────────────────────────────────────────
GREEN_DARK   = RGBColor(0x1B, 0x5E, 0x20)   # deep forest green
GREEN_MID    = RGBColor(0x2E, 0x7D, 0x32)   # medium green
GREEN_BRIGHT = RGBColor(0x66, 0xBB, 0x6A)   # bright green
TEAL         = RGBColor(0x00, 0x89, 0x7B)   # teal accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF1, 0xF8, 0xE9)   # very light green-white
DARK_TEXT    = RGBColor(0x1A, 0x23, 0x1A)   # near-black
YELLOW       = RGBColor(0xFF, 0xD6, 0x00)   # Samsung yellow accent
ORANGE       = RGBColor(0xFF, 0x6F, 0x00)   # warning orange

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper Functions ──────────────────────────────────────────────────────────

def bg(slide, color):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, alpha=None):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def txbox(slide, text, l, t, w, h,
          font_size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a transparent text box."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box


def heading(slide, text, top=Inches(0.3), size=36, color=WHITE):
    txbox(slide, text, Inches(0.5), top, Inches(12.3), Inches(0.7),
          font_size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def bullet_box(slide, lines, l, t, w, h,
               font_size=16, color=DARK_TEXT, bold_first=False):
    """Multi-line bullet text box."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold  = bold_first and (line == lines[0])


def divider(slide, top, color=GREEN_BRIGHT):
    rect(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), color)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title Slide
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)

# Top accent bar
rect(s, 0, 0, SLIDE_W, Inches(0.12), YELLOW)

# Bottom accent bar
rect(s, 0, Inches(7.38), SLIDE_W, Inches(0.12), YELLOW)

# Big background circle (decorative)
circle = s.shapes.add_shape(9, Inches(8.5), Inches(0.5), Inches(5), Inches(5))  # oval
circle.fill.solid()
circle.fill.fore_color.rgb = GREEN_MID
circle.line.fill.background()

# Title
txbox(s, "🌍  WASTE-NOT", Inches(0.8), Inches(1.4), Inches(11), Inches(1.3),
      font_size=54, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)

# Subtitle
txbox(s, "Circular Economy Tracker",
      Inches(0.8), Inches(2.7), Inches(9), Inches(0.8),
      font_size=30, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

divider(s, Inches(3.55), YELLOW)

txbox(s, "Samsung Hackathon  •  SDG 12: Responsible Consumption & Production",
      Inches(0.8), Inches(3.7), Inches(11), Inches(0.6),
      font_size=17, color=GREEN_BRIGHT, align=PP_ALIGN.LEFT, italic=True)

txbox(s, "Optimizing Waste Management Through\nData Analytics & Algorithms",
      Inches(0.8), Inches(4.4), Inches(9), Inches(1.0),
      font_size=20, color=WHITE, align=PP_ALIGN.LEFT)

# Bottom tag
txbox(s, "♻  Reduce  •  Recover  •  Reuse",
      Inches(0), Inches(6.7), Inches(13.33), Inches(0.5),
      font_size=14, color=YELLOW, align=PP_ALIGN.CENTER, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, SLIDE_W, Inches(1.1), GREEN_DARK)
heading(s, "🚨  The Problem We Are Solving", Inches(0.2), size=30)

# 3 problem cards
cards = [
    ("♻ Recycling\nContamination",
     "Up to 25% of recyclables\nare contaminated,\nrendering them useless\nand costly to process."),
    ("🚛 Inefficient\nLogistics",
     "Trucks run 40–60%\nempty on average,\nwasting fuel and\nincreasing CO₂ emissions."),
    ("📊 No Data-Driven\nDecisions",
     "Waste managers rely on\nguesswork instead of\nreal analytics to decide\ncollection routes & loads."),
]
for i, (title, body) in enumerate(cards):
    x = Inches(0.4 + i * 4.3)
    rect(s, x, Inches(1.3), Inches(4.0), Inches(5.5), GREEN_MID)
    txbox(s, title, x + Inches(0.15), Inches(1.5), Inches(3.7), Inches(1.1),
          font_size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    divider(s, Inches(2.65) + i*0, GREEN_BRIGHT)
    txbox(s, body, x + Inches(0.15), Inches(2.8), Inches(3.7), Inches(3.5),
          font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "Waste-Not solves all three with algorithms + data analytics.",
      Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
      font_size=15, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Our Solution Overview
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, SLIDE_W, Inches(1.1), TEAL)
heading(s, "💡  Our Solution: Waste-Not Tracker", Inches(0.2), size=30)

modules = [
    ("📥 1. Data Input", "Upload a CSV with site names, waste types, weight & contamination %.\nNo coding needed — just fill in the spreadsheet."),
    ("☣ 2. Contaminant\n    Detection", "Set operations (Union, Intersection) find materials present in 80%+ of sites — the common pollutants."),
    ("🔍 3. Unique Material\n    Analysis", "Set Difference identifies what each site uniquely produces for targeted collection strategies."),
    ("🚛 4. Greedy Truck\n    Loading", "Algorithm scores each batch by weight × cleanliness, fills the truck to maximum capacity efficiently."),
    ("📊 5. Venn Diagram\n    Visualization", "3-circle Venn diagram shows material overlap between Urban, Commercial & Industrial zones."),
]

for i, (title, body) in enumerate(modules):
    row = i // 3
    col = i  % 3
    x = Inches(0.35 + col * 4.33)
    y = Inches(1.25 + row * 2.8)
    rect(s, x, y, Inches(4.1), Inches(2.5), GREEN_MID)
    txbox(s, title, x+Inches(0.1), y+Inches(0.1), Inches(3.9), Inches(0.8),
          font_size=16, bold=True, color=YELLOW)
    txbox(s, body, x+Inches(0.1), y+Inches(0.85), Inches(3.9), Inches(1.5),
          font_size=13, color=WHITE)

txbox(s, "All modules feed from one CSV → automated reports + charts → actionable decisions.",
      Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35),
      font_size=13, color=GREEN_BRIGHT, align=PP_ALIGN.CENTER, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – CSV Input Format
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, SLIDE_W, Inches(1.1), GREEN_DARK)
heading(s, "📥  Data Input: waste_input_data.csv", Inches(0.2), size=28)

# CSV column explanation
cols = [
    ("site_id",           "A, B, C …",                       "Short unique identifier"),
    ("site_name",         "Site_A, Site_B …",                "Human-readable name"),
    ("site_type",         "Urban / Industrial / Commercial",  "Zone classification"),
    ("materials",         "plastic_bottles|cardboard|glass",  "Pipe-separated waste types"),
    ("total_weight_kg",   "450, 680, 750 …",                 "Total waste weight in kg"),
    ("contamination_pct", "15, 22, 8 …",                     "% of incorrectly sorted waste"),
]

rect(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.5), GREEN_MID)
for ci, hdr in enumerate(["Column Name", "Example Value", "What It Means"]):
    txbox(s, hdr, Inches(0.4 + ci*4.2), Inches(1.25), Inches(4.0), Inches(0.4),
          font_size=15, bold=True, color=YELLOW)

for ri, (col, ex, desc) in enumerate(cols):
    bg_c = GREEN_DARK if ri % 2 == 0 else GREEN_MID
    rect(s, Inches(0.3), Inches(1.75 + ri*0.75), Inches(12.7), Inches(0.72), bg_c)
    txbox(s, col,  Inches(0.4),        Inches(1.79 + ri*0.75), Inches(4.0), Inches(0.65), font_size=14, bold=True,  color=GREEN_BRIGHT)
    txbox(s, ex,   Inches(4.6),        Inches(1.79 + ri*0.75), Inches(4.0), Inches(0.65), font_size=13, color=WHITE)
    txbox(s, desc, Inches(8.7),        Inches(1.79 + ri*0.75), Inches(4.2), Inches(0.65), font_size=13, color=LIGHT_GRAY)

txbox(s, "💡  Just edit the CSV and re-run — no code changes needed!",
      Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.45),
      font_size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Contaminant Detection
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, SLIDE_W, Inches(1.1), ORANGE)
heading(s, "☣  Contaminant Identification Analysis", Inches(0.2), size=28, color=WHITE)

# Left: explanation box
rect(s, Inches(0.3), Inches(1.2), Inches(5.5), Inches(5.8), GREEN_MID)
txbox(s, "How It Works", Inches(0.45), Inches(1.3), Inches(5.2), Inches(0.5),
      font_size=18, bold=True, color=YELLOW)
steps = [
    "1. Collect all material sets from every site",
    "2. Union → find all unique materials across sites",
    "3. Count how many sites each material appears in",
    "4. Threshold: appears in 4+ sites (≥80%) = CONTAMINANT",
    "5. Intersection → materials in ALL 5 sites flagged",
]
bullet_box(s, steps, Inches(0.45), Inches(1.9), Inches(5.2), Inches(3.0),
           font_size=14, color=WHITE)

txbox(s, "Python operation:\nset.intersection(*all_sets)\nset.union(*all_sets)",
      Inches(0.45), Inches(5.0), Inches(5.2), Inches(0.9),
      font_size=13, color=GREEN_BRIGHT, italic=True)

# Right: results table
rect(s, Inches(6.1), Inches(1.2), Inches(7.0), Inches(5.8), RGBColor(0x1A,0x23,0x1A))
txbox(s, "Results from our data:", Inches(6.25), Inches(1.3), Inches(6.7), Inches(0.45),
      font_size=16, bold=True, color=YELLOW)

results = [
    ("cardboard",       "5/5  100%", "⚠ CONTAMINANT", ORANGE),
    ("plastic_bottles", "5/5  100%", "⚠ CONTAMINANT", ORANGE),
    ("plastic_bags",    "4/5   80%", "⚠ CONTAMINANT", ORANGE),
    ("aluminum_cans",   "4/5   80%", "⚠ CONTAMINANT", ORANGE),
    ("glass",           "4/5   80%", "⚠ CONTAMINANT", ORANGE),
    ("electronics",     "4/5   80%", "⚠ CONTAMINANT", ORANGE),
    ("paper",           "4/5   80%", "⚠ CONTAMINANT", ORANGE),
    ("textiles",        "3/5   60%", "✅ Normal",      GREEN_BRIGHT),
    ("food_waste",      "2/5   40%", "✅ Normal",      GREEN_BRIGHT),
]
for ri, (mat, freq, status, sc) in enumerate(results):
    y = Inches(1.82 + ri * 0.52)
    txbox(s, mat,    Inches(6.2),  y, Inches(3.0), Inches(0.48), font_size=13, color=WHITE)
    txbox(s, freq,   Inches(9.3),  y, Inches(1.5), Inches(0.48), font_size=13, color=WHITE)
    txbox(s, status, Inches(10.8), y, Inches(2.2), Inches(0.48), font_size=12, bold=True, color=sc)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Greedy Truck Loading
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, SLIDE_W, Inches(1.1), TEAL)
heading(s, "🚛  Greedy Truck Loading Algorithm", Inches(0.2), size=28)

# Formula box
rect(s, Inches(0.3), Inches(1.2), Inches(6.2), Inches(2.2), GREEN_DARK)
txbox(s, "Priority Formula", Inches(0.45), Inches(1.28), Inches(5.9), Inches(0.5),
      font_size=17, bold=True, color=YELLOW)
txbox(s, "Priority  =  Weight  ×  (1 − contamination%)",
      Inches(0.45), Inches(1.82), Inches(5.9), Inches(0.6),
      font_size=15, bold=True, color=GREEN_BRIGHT)
txbox(s, "Higher weight + Lower contamination = Higher priority\n(Cleanest & heaviest goes first)",
      Inches(0.45), Inches(2.45), Inches(5.9), Inches(0.75),
      font_size=13, color=WHITE)

# stats box
rect(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.2), GREEN_DARK)
stats = [
    "🏋 Total waste available : 2,780 kg",
    "🚛 Truck capacity           : 1,000 kg",
    "📦 Total batches created : 12",
    "✅ Batches loaded            : 4",
    "📈 Utilization                    : 99%",
]
bullet_box(s, stats, Inches(6.95), Inches(1.28), Inches(5.9), Inches(2.1),
           font_size=14, color=WHITE)

# Loading steps table
rect(s, Inches(0.3), Inches(3.6), Inches(12.7), Inches(0.45), GREEN_MID)
for ci, hdr in enumerate(["Batch", "Weight", "Priority", "Running Total", "Decision"]):
    txbox(s, hdr, Inches(0.4 + ci*2.5), Inches(3.65), Inches(2.4), Inches(0.38),
          font_size=13, bold=True, color=YELLOW)

rows = [
    ("E_batch_1", "290.00 kg", "266.80", "290.00 kg",  "✅ LOADED",  GREEN_BRIGHT),
    ("E_batch_2", "290.00 kg", "266.80", "580.00 kg",  "✅ LOADED",  GREEN_BRIGHT),
    ("C_batch_1", "250.00 kg", "205.00", "830.00 kg",  "✅ LOADED",  GREEN_BRIGHT),
    ("C_batch_2", "250.00 kg", "205.00", "1080 kg",   "❌ SKIP",    ORANGE),
    ("A_batch_1", "225.00 kg", "191.25", "1055 kg",   "❌ SKIP",    ORANGE),
    ("B_batch_1", "226.67 kg", "176.80", "1056 kg",   "❌ SKIP",    ORANGE),
    ("D_batch_1", "160.00 kg", "140.80", "990.00 kg", "✅ LOADED",  GREEN_BRIGHT),
]
for ri, (b, w, p, rt, dec, dc) in enumerate(rows):
    bg_c = RGBColor(0x1A,0x23,0x1A) if ri%2==0 else GREEN_DARK
    rect(s, Inches(0.3), Inches(4.1+ri*0.41), Inches(12.7), Inches(0.4), bg_c)
    for ci, val in enumerate([b, w, p, rt]):
        txbox(s, val, Inches(0.4+ci*2.5), Inches(4.13+ri*0.41), Inches(2.4), Inches(0.37),
              font_size=12, color=WHITE)
    txbox(s, dec, Inches(10.4), Inches(4.13+ri*0.41), Inches(2.5), Inches(0.37),
          font_size=12, bold=True, color=dc)

txbox(s, "Final result: 990 kg / 1000 kg  =  99% truck utilization  🎯",
      Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.4),
      font_size=15, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Venn Diagram Slide
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, SLIDE_W, Inches(1.1), GREEN_MID)
heading(s, "📊  Material Overlap: Venn Diagram Analysis", Inches(0.2), size=28)

# Insert the venn diagram image if it exists
img_path = os.path.join(os.path.dirname(__file__), 'waste_materials_venn_diagram.png')
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.15), Inches(7.5), Inches(5.5))
else:
    txbox(s, "[Venn diagram image not found.\nRun waste_not.py first to generate it.]",
          Inches(0.3), Inches(2.0), Inches(7.5), Inches(2.0),
          font_size=16, color=ORANGE, align=PP_ALIGN.CENTER)

# Overlap stats on right
rect(s, Inches(8.1), Inches(1.2), Inches(5.0), Inches(5.5), RGBColor(0x1A,0x23,0x1A))
txbox(s, "Overlap Statistics", Inches(8.25), Inches(1.3), Inches(4.7), Inches(0.5),
      font_size=17, bold=True, color=YELLOW)

stats2 = [
    ("A ∩ B  (Urban + Commercial)",  "7 materials"),
    ("A ∩ C  (Urban + Industrial)",  "4 materials"),
    ("B ∩ C  (Commercial + Industrial)", "5 materials"),
    ("A ∩ B ∩ C  (ALL THREE)",       "4 materials ⭐"),
]
for ri, (label, val) in enumerate(stats2):
    y = Inches(1.9 + ri*0.9)
    rect(s, Inches(8.2), y, Inches(4.8), Inches(0.8), GREEN_MID)
    txbox(s, label, Inches(8.3), y+Inches(0.03), Inches(3.0), Inches(0.38),
          font_size=12, color=WHITE)
    txbox(s, val,   Inches(11.1), y+Inches(0.03), Inches(1.8), Inches(0.38),
          font_size=13, bold=True, color=GREEN_BRIGHT)

txbox(s, "Core universal materials (in all 3 sites):\ncardboard, plastic_bottles, electronics, plastic_bags",
      Inches(8.2), Inches(5.6), Inches(4.9), Inches(0.9),
      font_size=13, color=YELLOW, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Industry Impact
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, SLIDE_W, Inches(1.1), GREEN_DARK)
heading(s, "🏭  Real-World Industry Impact", Inches(0.2), size=28)

impacts = [
    ("💰 Fuel\nSavings",
     "99% truck utilization\nvs industry avg 60%.\nSaves ~33% fuel costs\nper fleet per day."),
    ("🌱 CO₂\nReduction",
     "Fewer trips = fewer\nemissions. Supports\ncompany carbon\nneutrality targets."),
    ("♻ Higher\nRevenue",
     "Clean-first loading\nmeans recyclers get\nbetter material fetching\n30–50% higher prices."),
    ("🎓 Targeted\nEducation",
     "Contaminant data\nshows WHICH materials\nneed public awareness\ncampaigns."),
    ("🗺 Smart\nRouting",
     "Unique material maps\nshow which sites need\nspecialized vehicles\n(hazardous, e-waste)."),
    ("📦 Logistics\nAutomation",
     "Replaces manual\nguesswork with data-\ndriven loading plans\nevery single day."),
]
for i, (title, body) in enumerate(impacts):
    row = i // 3
    col = i  % 3
    x = Inches(0.35 + col * 4.33)
    y = Inches(1.25 + row * 2.85)
    rect(s, x, y, Inches(4.1), Inches(2.6), GREEN_DARK)
    txbox(s, title, x+Inches(0.12), y+Inches(0.1), Inches(3.86), Inches(0.75),
          font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    divider(s, y + Inches(0.88), TEAL)
    txbox(s, body, x+Inches(0.12), y+Inches(0.98), Inches(3.86), Inches(1.5),
          font_size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Scalability & Future Scope
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, SLIDE_W, Inches(1.1), TEAL)
heading(s, "🚀  Scalability & Future Scope", Inches(0.2), size=28)

# Left column: Current vs Scale
rect(s, Inches(0.3), Inches(1.2), Inches(5.9), Inches(5.8), GREEN_MID)
txbox(s, "Current Demo", Inches(0.45), Inches(1.3), Inches(5.6), Inches(0.5),
      font_size=17, bold=True, color=YELLOW)
current = [
    "📄 5 sites, manual CSV input",
    "🚛 1 truck, fixed 1000 kg capacity",
    "📊 Static Venn diagram output",
    "🖥 Run manually via terminal",
    "📁 3 output files generated",
]
bullet_box(s, current, Inches(0.45), Inches(1.9), Inches(5.6), Inches(2.5),
           font_size=14, color=WHITE)

txbox(s, "Future Scale", Inches(0.45), Inches(4.5), Inches(5.6), Inches(0.5),
      font_size=17, bold=True, color=YELLOW)
future = [
    "📡 500+ sites with IoT sensors (auto weight)",
    "🚛 50+ trucks, dynamic capacities per type",
    "📊 Real-time dashboard (web-based)",
    "⚡ Scheduled daily automated runs",
    "🤖 ML-based contamination prediction",
]
bullet_box(s, future, Inches(0.45), Inches(5.1), Inches(5.6), Inches(2.1),
           font_size=14, color=GREEN_BRIGHT)

# Right column: Tech stack / roadmap
rect(s, Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.8), RGBColor(0x1A,0x23,0x1A))
txbox(s, "Technology Roadmap", Inches(6.65), Inches(1.3), Inches(6.2), Inches(0.5),
      font_size=17, bold=True, color=YELLOW)
phases = [
    ("Phase 1 (Now)",      "Python + CSV  → reports & charts"),
    ("Phase 2",            "REST API + Database → multi-user access"),
    ("Phase 3",            "IoT Integration → real-time weight feeds"),
    ("Phase 4",            "Web Dashboard → live analytics portal"),
    ("Phase 5",            "ML Models → predict contamination risk"),
]
for ri, (phase, desc) in enumerate(phases):
    rect(s, Inches(6.6), Inches(1.9 + ri*1.0), Inches(6.2), Inches(0.9),
         GREEN_MID if ri % 2 == 0 else GREEN_DARK)
    txbox(s, phase, Inches(6.7), Inches(1.95 + ri*1.0), Inches(2.0), Inches(0.38),
          font_size=13, bold=True, color=YELLOW)
    txbox(s, desc,  Inches(8.8), Inches(1.95 + ri*1.0), Inches(4.0), Inches(0.38),
          font_size=13, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Thank You / Closing
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, SLIDE_W, Inches(0.12), YELLOW)
rect(s, 0, Inches(7.38), SLIDE_W, Inches(0.12), YELLOW)

# Large circle decoration
circle2 = s.shapes.add_shape(9, Inches(-1), Inches(4.5), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = GREEN_MID
circle2.line.fill.background()

txbox(s, "🌍", Inches(4.0), Inches(0.8), Inches(5), Inches(1.2),
      font_size=60, align=PP_ALIGN.CENTER, color=WHITE)

txbox(s, "Thank You!", Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.1),
      font_size=52, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

txbox(s, "WASTE-NOT: Circular Economy Tracker",
      Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.7),
      font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

divider(s, Inches(3.9), YELLOW)

summary = [
    "✅  Detects contaminants using Set Operations",
    "✅  Optimizes truck loading with Greedy Algorithm",
    "✅  Visualizes material flow with Venn Diagrams",
    "✅  Fully data-driven — just update the CSV!",
]
bullet_box(s, summary, Inches(3.0), Inches(4.1), Inches(7.3), Inches(2.4),
           font_size=17, color=WHITE)

txbox(s, "Supporting SDG 12 · Reducing Waste · Saving the Planet ♻",
      Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
      font_size=14, color=GREEN_BRIGHT, align=PP_ALIGN.CENTER, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), 'WasteNot_Presentation.pptx')
prs.save(out_path)
print(f"✅ Presentation saved: {out_path}")
print(f"   Slides created: {len(prs.slides)}")
