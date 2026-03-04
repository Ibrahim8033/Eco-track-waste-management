"""
add_slides.py  –  Adds 3 new premium slides to WasteNot_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── palette ──────────────────────────────────────────────────────────────────
G_DARK   = RGBColor(0x0D,0x2B,0x0D)
G_MID    = RGBColor(0x1B,0x5E,0x20)
G_LIGHT  = RGBColor(0x2E,0x7D,0x32)
G_BRIGHT = RGBColor(0x66,0xBB,0x6A)
TEAL     = RGBColor(0x00,0x89,0x7B)
YELLOW   = RGBColor(0xFF,0xD6,0x00)
ORANGE   = RGBColor(0xFF,0x6F,0x00)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
LIGHT_BG = RGBColor(0xF1,0xF8,0xE9)
BLUE_ACC = RGBColor(0x29,0xB6,0xF6)
PURPLE   = RGBColor(0x7C,0x4D,0xFF)
PINK     = RGBColor(0xEC,0x40,0x7A)

# ── helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line_color=None, shape_type=1):
    sh = slide.shapes.add_shape(shape_type, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def hdr(slide, text, color=WHITE, size=28):
    box(slide, 0, 0, Inches(13.33), Inches(1.05), G_MID)
    # accent left bar
    box(slide, 0, 0, Inches(0.18), Inches(1.05), YELLOW)
    txt(slide, text, Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.75),
        size=size, bold=True, color=color, align=PP_ALIGN.LEFT)

def divline(slide, t, c=G_BRIGHT):
    box(slide, Inches(0.4), t, Inches(12.53), Inches(0.04), c)

# ── load existing pptx ────────────────────────────────────────────────────────
path = r"c:\Users\Admin\Desktop\hackathon samsung\temp1\WasteNot_Presentation.pptx"
prs  = Presentation(path)
BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# NEW SLIDE A  –  "By The Numbers" (Stats Dashboard)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, G_DARK)

# Full-width gradient-style top bar
box(s, 0, 0, Inches(13.33), Inches(1.05), G_MID)
box(s, 0, 0, Inches(0.18),  Inches(1.05), YELLOW)
box(s, 0, Inches(1.05), Inches(13.33), Inches(0.06), YELLOW)
txt(s, "📈  KEY NUMBERS AT A GLANCE", Inches(0.3), Inches(0.18),
    Inches(12), Inches(0.72), size=26, bold=True, color=YELLOW)

# big stat cards — row 1
cards1 = [
    (TEAL,   "5",       "Collection Sites\nAnalyzed"),
    (PURPLE, "15",      "Total Material\nTypes Tracked"),
    (ORANGE, "7",       "Common\nContaminants\nDetected"),
    (PINK,   "100%",    "Cardboard &\nPlastic Bottles\nUniversal Rate"),
]
for i,(c,num,label) in enumerate(cards1):
    x = Inches(0.28 + i*3.27)
    box(s, x, Inches(1.22), Inches(3.1), Inches(2.5), c)
    box(s, x, Inches(1.22), Inches(3.1), Inches(0.08), YELLOW)   # top accent
    txt(s, num,   x+Inches(0.1), Inches(1.35), Inches(2.9), Inches(1.05),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, x+Inches(0.1), Inches(2.42), Inches(2.9), Inches(0.88),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)

# big stat cards — row 2
cards2 = [
    (G_LIGHT,            "2,780 kg", "Total Waste\nAvailable"),
    (RGBColor(0x00,0x7A,0x7A), "12",      "Batches\nCreated"),
    (RGBColor(0xD8,0x4A,0x00), "99%",     "Truck\nUtilization"),
    (RGBColor(0x4A,0x00,0x8C), "4",       "Optimal Batches\nSelected"),
]
for i,(c,num,label) in enumerate(cards2):
    x = Inches(0.28 + i*3.27)
    box(s, x, Inches(3.90), Inches(3.1), Inches(2.5), c)
    box(s, x, Inches(3.90), Inches(3.1), Inches(0.08), WHITE)
    txt(s, num,   x+Inches(0.1), Inches(4.03), Inches(2.9), Inches(1.05),
        size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, x+Inches(0.1), Inches(5.1),  Inches(2.9), Inches(0.88),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Waste-Not turns raw waste data into these measurable outcomes  🌱",
    Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.5),
    size=14, italic=True, color=G_BRIGHT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# NEW SLIDE B  –  Set Operations Visual (Math Behind The Magic)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
hdr(s, "🔢  Set Operations: The Math Behind the Magic", color=WHITE, size=26)

ops = [
    (TEAL,   "UNION  ( A ∪ B )",
     "Combines ALL materials from\nboth sites into one big set.\n\n"
     "→ Gives us the complete list\n   of every material collected\n   across the network.",
     "A = {plastic, glass, paper}\nB = {plastic, metal, glass}\n\n"
     "A ∪ B = {plastic, glass,\n             paper, metal}"),

    (PURPLE, "INTERSECTION  ( A ∩ B )",
     "Finds materials present in\nBOTH sites simultaneously.\n\n"
     "→ Used to detect CONTAMINANTS\n   that appear in 80%+ sites.",
     "A = {plastic, glass, paper}\nB = {plastic, metal, glass}\n\n"
     "A ∩ B = {plastic, glass}"),

    (ORANGE, "DIFFERENCE  ( A − B )",
     "Finds materials UNIQUE to\none site only.\n\n"
     "→ Identifies site-specific waste\n   for targeted collection strategy.",
     "A = {plastic, glass, paper}\nB = {plastic, metal, glass}\n\n"
     "A − B = {paper}"),
]

for i,(c,title,expl,example) in enumerate(ops):
    x = Inches(0.25 + i*4.36)
    # main card
    box(s, x, Inches(1.15), Inches(4.15), Inches(5.8), c)
    # top accent
    box(s, x, Inches(1.15), Inches(4.15), Inches(0.1), YELLOW)
    # title
    txt(s, title, x+Inches(0.12), Inches(1.3), Inches(3.9), Inches(0.65),
        size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    divline(s, Inches(2.02) if i==0 else Inches(2.02), c)
    # explanation
    txt(s, expl, x+Inches(0.12), Inches(2.08), Inches(3.9), Inches(2.3),
        size=13, color=WHITE)
    # example box inside card
    box(s, x+Inches(0.12), Inches(4.45), Inches(3.9), Inches(2.1), G_DARK)
    txt(s, "📦 Example:", x+Inches(0.22), Inches(4.5), Inches(3.7), Inches(0.38),
        size=12, bold=True, color=YELLOW)
    txt(s, example, x+Inches(0.22), Inches(4.9), Inches(3.7), Inches(1.55),
        size=12, color=G_BRIGHT, italic=True)

txt(s, "These 3 operations power the entire contaminant detection engine of Waste-Not",
    Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.38),
    size=13, bold=True, color=G_MID, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# NEW SLIDE C  –  How It Works Pipeline (Flow Diagram)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, G_DARK)
hdr(s, "⚙  How It Works: End-to-End Pipeline", color=WHITE, size=26)
box(s, 0, Inches(1.05), Inches(13.33), Inches(0.06), YELLOW)

steps = [
    (TEAL,   "📥",  "INPUT",         "Upload\nwaste_input_data.csv\nwith site data"),
    (BLUE_ACC,"🔍", "DETECT",        "Set operations\nfind contaminants\nin 80%+ sites"),
    (PURPLE, "📐",  "ANALYSE",       "Set difference\nfinds unique\nmaterials per site"),
    (ORANGE, "🚛",  "OPTIMISE",      "Greedy algorithm\nloads truck at\n99% utilization"),
    (PINK,   "📊",  "VISUALISE",     "Venn diagram\nshows material\noverlaps"),
    (G_BRIGHT,"📄", "REPORT",        "3 output files:\nreport, manifest,\ndiagram"),
]

arrow_y = Inches(4.05)
for i,(c,icon,step,desc) in enumerate(steps):
    x = Inches(0.22 + i*2.19)
    # vertical connector lines (not for last)
    if i < len(steps)-1:
        box(s, x+Inches(2.05), arrow_y+Inches(0.18),
            Inches(0.14), Inches(0.3), YELLOW)
        # arrowhead triangle approximation
        box(s, x+Inches(2.09), arrow_y+Inches(0.46),
            Inches(0.06), Inches(0.12), YELLOW)

    # circle icon
    circ = s.shapes.add_shape(9, x, Inches(1.25), Inches(2.0), Inches(2.0))
    circ.fill.solid(); circ.fill.fore_color.rgb = c
    circ.line.fill.background()
    txt(s, icon, x, Inches(1.55), Inches(2.0), Inches(0.8),
        size=26, align=PP_ALIGN.CENTER, color=WHITE)

    # step number badge
    badge = s.shapes.add_shape(9, x+Inches(1.5), Inches(1.15),
                                Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = YELLOW
    badge.line.fill.background()
    txt(s, str(i+1), x+Inches(1.5), Inches(1.18), Inches(0.45), Inches(0.38),
        size=13, bold=True, color=G_DARK, align=PP_ALIGN.CENTER)

    # step label
    txt(s, step, x, Inches(3.32), Inches(2.0), Inches(0.45),
        size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # desc card
    box(s, x, Inches(3.82), Inches(2.0), Inches(2.35), G_MID)
    txt(s, desc, x+Inches(0.08), Inches(3.92), Inches(1.85), Inches(2.1),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)

# bottom data flow line
box(s, Inches(0.22), Inches(6.35), Inches(13.0), Inches(0.06), G_BRIGHT)
txt(s, "⬅─────────────────────  DATA FLOWS LEFT TO RIGHT  ─────────────────────➡",
    Inches(0.2), Inches(6.45), Inches(12.9), Inches(0.4),
    size=12, italic=True, color=G_BRIGHT, align=PP_ALIGN.CENTER)

txt(s, "Fully automated · Just update the CSV and all 6 stages run instantly",
    Inches(0.2), Inches(6.92), Inches(12.9), Inches(0.42),
    size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\Admin\Desktop\hackathon samsung\temp1\WasteNot_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Total slides now: {len(prs.slides)}")
